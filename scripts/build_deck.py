"""
Build the pitch deck from a COPY of the mandatory Redrob template.

The template (`Idea Submission Template _ Redrob.pptx`) is treated as read-only:
we copy it to deck/india-runs-hackathon-deck.pptx and only fill the existing
content placeholders. Titles, branding images, layout, fonts and colours are
left exactly as the template defines them, per the rule that the template must
stay constant.

Team-identity and URL fields are intentionally left as <PLACEHOLDERS> for the
participant to swap in before submission.

    python scripts/build_deck.py
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Idea Submission Template _ Redrob.pptx"
OUT = ROOT / "deck" / "india-runs-hackathon-deck.pptx"

# Match the template's body styling so filled text blends with the design.
BODY_FONT = "Manrope SemiBold"
MONO_FONT = "Consolas"
INK = RGBColor(0x20, 0x27, 0x29)
ACCENT = RGBColor(0x6D, 0x4A, 0xFF)

# Placeholders the participant fills in before submitting.
TEAM_NAME = "<TEAM_NAME = your Team ID>"
TEAM_LEADER = "<TEAM_LEADER_NAME>"
GITHUB_URL = "https://github.com/harshkawatra11/india-runs-hackathon"
VERCEL_URL = "https://india-runs-hackathon.vercel.app"
VIDEO_URL = "<demo video URL>"


def _no_auto_bullet(para):
    """Strip any inherited auto-bullet so we fully control the glyphs."""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def set_body(shape, bullets, *, size=11, mono=False, accent_first=False):
    """Replace a text frame's content with cleanly, uniformly styled bullets.

    `bullets` is a list of (text, indent_level) tuples or plain strings. We
    render our own bullet glyphs (level 0 -> round dot, level 1 -> en-dash with
    indent) and disable the placeholder's inherited auto-bullets so spacing and
    markers are consistent across every line.
    """
    tf = shape.text_frame
    tf.word_wrap = True
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    for i, item in enumerate(bullets):
        text, level = (item if isinstance(item, tuple) else (item, 0))
        para = first if i == 0 else tf.add_paragraph()
        para.level = level
        para.line_spacing = 0.95 if mono else 1.05
        para.space_after = Pt(1) if mono else Pt(5)
        para.space_before = Pt(0)
        if not mono:
            _no_auto_bullet(para)
            glyph = "      –  " if level >= 1 else "•  "
            text = glyph + text
        run = para.add_run()
        run.text = text
        f = run.font
        f.name = MONO_FONT if mono else BODY_FONT
        f.size = Pt(size)
        f.color.rgb = ACCENT if (accent_first and i == 0) else INK


def append_inline(shape, suffix):
    """Append text to the end of an existing single-line box (slide 1 fields)."""
    para = shape.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = "  " + suffix
    # copy styling from the existing label run if present
    if para.runs:
        src = para.runs[0].font
        run.font.name = src.name or BODY_FONT
        if src.size:
            run.font.size = src.size
        run.font.color.rgb = INK


def by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(shape_id)


def add_textbox(slide, left_in, top_in, width_in, height_in, bullets, **kw):
    from pptx.util import Inches
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                   Inches(width_in), Inches(height_in))
    set_body(box, bullets, **kw)
    return box


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(TEMPLATE, OUT)
    prs = Presentation(str(OUT))
    slides = prs.slides

    # ---- Slide 1: identity ------------------------------------------------
    s1 = slides[0]
    append_inline(by_id(s1, 55), TEAM_NAME)              # Team Name :
    append_inline(by_id(s1, 57), TEAM_LEADER)            # Team Leader Name :
    append_inline(by_id(s1, 56),                         # Problem Statement :
                  "Intelligent Candidate Discovery & Ranking — rank the top-100 "
                  "best-fit candidates for a nuanced Senior AI Engineer JD from a "
                  "100,000-candidate pool, going far beyond keyword matching.")

    # ---- Slide 2: Solution Overview --------------------------------------
    set_body(by_id(slides[1], 64), [
        "A deterministic, fully-explainable ranking engine that encodes the JD as an auditable rubric and scores every candidate on semantic fit + structured evidence, then modulates that score by real-world availability and a honeypot-integrity check.",
        "Runs CPU-only with no network or LLM calls at rank time: 100,000 candidates ranked in ~98 seconds — comfortably inside the 5-min / 16 GB budget.",
        "What differentiates it from keyword/embedding matching:",
        ("Title is a GATE, not a feature — an “HR Manager” with 9 AI skills cannot outrank a real ML engineer.", 1),
        ("Skills are trusted only when endorsed AND actually used — defusing keyword stuffing.", 1),
        ("Impossible “honeypot” profiles collapse via general consistency checks, not special-casing.", 1),
        ("Behavioural signals modulate the score, so a brilliant-but-inactive profile is correctly down-weighted.", 1),
    ], size=11)

    # ---- Slide 3: JD Understanding & Candidate Evaluation ----------------
    set_body(by_id(slides[2], 71), [
        "Key requirements extracted from the JD:",
        ("5–9 yrs (ideal 6–8); production embeddings-retrieval + vector DB + ranking-evaluation (NDCG/MRR/MAP); strong Python; shipped real search/ranking/recsys at PRODUCT companies; Pune/Noida or Tier-1 India / open to relocation; genuinely available.", 1),
        "Explicit hard-negatives the JD names (and we encode):",
        ("keyword-stuffers in non-eng roles, pure-research-only, LangChain-only juniors, title-chasers, consulting-only careers (TCS/Infosys/…), CV/speech/robotics without NLP/IR.", 1),
        "Most important signals for relevance:",
        ("current title + career evidence of systems actually built, trusted-skill match, company type, and availability (last-active, recruiter response rate, open-to-work).", 1),
        "Beyond keywords: we read the gap between what the JD says and what it means — a plain-language profile that built a recommender at a product company outranks a buzzword-stuffed one.",
    ], size=10.5)

    # ---- Slide 4: Ranking Methodology ------------------------------------
    set_body(by_id(slides[3], 78), [
        "Represent & retrieve: concatenate profile + career text; measure semantic similarity to the JD via a hybrid of lexical TF-IDF, latent-semantic LSA (TruncatedSVD), and an OPTIONAL precomputed MiniLM dense embedding — all computed offline.",
        "Score (weighted base, sums to 1.0):",
        ("semantic 0.26 · role-fit 0.24 · skill-trust 0.20 · experience 0.12 · company 0.10 · location 0.05 · education 0.03.", 1),
        "Modulate multiplicatively:",
        ("final = base × disqualifier-penalty × availability-multiplier × integrity-multiplier — so any single fatal flaw (off-track title, honeypot) collapses an otherwise glossy profile.", 1),
        "Rank: sort by score descending, break ties by candidate_id ascending, emit a spec-exact top-100 CSV and self-validate it.",
    ], size=10.5)

    # ---- Slide 5: Explainability & Data Validation -----------------------
    set_body(by_id(slides[4], 85), [
        "Every rank ships a 1–2 sentence reasoning assembled ONLY from the candidate’s real facts — years, current title, named matched skills, a concrete signal value — varied per row and rank-consistent in tone.",
        "No hallucination by construction: reasoning is generated from extracted fields (no LLM), so every claim provably traces to the profile — directly satisfying the Stage-4 manual-review checks.",
        "Suspicious / low-quality / honeypot profiles: the integrity module flags impossible date spans, experience-math contradictions and “expert with 0 months used”, collapsing the score. Result: 0 honeypots in our top-100.",
        "Keyword stuffers defused: a non-engineering current title gates role-fit, and untrusted skills (no endorsements / no usage) contribute almost nothing.",
    ], size=10.5)

    # ---- Slide 6: End-to-End Workflow ------------------------------------
    set_body(by_id(slides[5], 92), [
        "1.  Stream candidates.jsonl line-by-line (constant memory over the 487 MB file).",
        "2.  Extract per-candidate features + integrity check + availability modifier.",
        "3.  Vectorise all profile text; compute JD semantic similarity (TF-IDF + LSA [+ optional dense]).",
        "4.  Compose the final score (base × penalty × availability × integrity).",
        "5.  Rank, break ties by candidate_id, take the top 100.",
        "6.  Generate fact-grounded reasoning for each survivor.",
        "7.  Write a spec-valid CSV and run the official validator.",
        "One command:  python -m engine.rank --candidates candidates.jsonl --out submission.csv",
    ], size=11)

    # ---- Slide 7: System Architecture (diagram) --------------------------
    diagram = [
        "candidates.jsonl  (100,000 profiles, streamed)",
        "        |",
        "        v",
        "  FEATURE EXTRACTION ---+-- role / title gate",
        "                        +-- skill-trust  (endorsements x duration x assessment)",
        "                        +-- company type / location / experience band",
        "  INTEGRITY CHECK ------+-- honeypot consistency (dates, exp-math, 0-use)",
        "  AVAILABILITY ---------+-- recency, response rate, open-to-work, notice",
        "        |",
        "        v",
        "  SEMANTIC ENGINE :  TF-IDF  +  LSA   ( + optional MiniLM dense, precomputed )",
        "        |",
        "        v",
        "  SCORE COMPOSITION :  base x penalty x availability x integrity",
        "        |",
        "        v",
        "  RANK + tie-break  ->  top-100  ->  grounded reasoning  ->  submission.csv",
        "                        |",
        "                        +--->  Next.js + Vercel sandbox  (same logic, ported to TypeScript)",
    ]
    add_textbox(slides[6], 0.41, 1.45, 9.3, 3.6,
                [(line, 0) for line in diagram], size=9, mono=True)

    # ---- Slide 8: Results & Performance ----------------------------------
    set_body(by_id(slides[7], 105), [
        "Ranking quality: the top-100 are senior AI / ML / Search / NLP engineers at product companies (Paytm, Apple, Microsoft, Netflix, Google, Meta) and AI startups — 6–8 yrs, real retrieval/ranking evidence, strong availability.",
        "Trap resistance: 0 honeypots in the top-100; the sample’s planted “HR Manager + 9 AI skills” stuffer is correctly buried.",
        "Compute: 100,000 candidates ranked in ~85 s on CPU with the full hybrid (auto-loading a precomputed dense cache; ~98 s on the lexical+LSA fallback), peak RAM well under 16 GB — inside the 5-min / 16 GB / no-GPU / no-network budget.",
        "Metric-aware: top-heavy precision tuned for NDCG@10 / NDCG@50; the official validator passes (exactly 100 rows, unique ranks 1–100, non-increasing score, id tie-break).",
    ], size=10.5)

    # ---- Slide 9: Technologies Used --------------------------------------
    set_body(by_id(slides[8], 112), [
        "Python 3 + NumPy + scikit-learn (TF-IDF, TruncatedSVD/LSA): fast, deterministic, CPU-only, zero heavy dependencies at rank time. Chosen because the JD itself prizes retrieval & evaluation fundamentals over framework hype.",
        "sentence-transformers (MiniLM): OPTIONAL offline dense embeddings, precomputed to a cached .npy so the rank step stays network-free.",
        "pytest: behavioural test suite (honeypot detection, keyword-stuffer rejection, tie-break ordering, CSV validity against the official validator).",
        "Next.js 15 + React 19 + TypeScript + Tailwind on Vercel: the recruiter-facing showcase and reproducibility sandbox, running the same scoring logic ported to TypeScript.",
        "No hosted-LLM calls anywhere in the ranking path — both for spec compliance and production realism (a system that calls an LLM per candidate cannot scale).",
    ], size=10.5)

    # ---- Slide 10: Submission Assets -------------------------------------
    set_body(by_id(slides[9], 119), [
        ("GitHub repository:  " + GITHUB_URL, 0),
        ("Live demo / sandbox (Vercel):  " + VERCEL_URL, 0),
        ("Reproduce command:  python -m engine.rank --candidates candidates.jsonl --out submission.csv   (< 5 min, CPU)", 0),
        ("Artifacts:  submission.csv (top-100) · submission_metadata.yaml · requirements.txt · engine/ + tests/", 0),
        ("Demo video:  " + VIDEO_URL, 0),
    ], size=11)

    prs.save(str(OUT))
    print(f"Deck written to {OUT}  ({len(slides)} slides, template design preserved)")


if __name__ == "__main__":
    main()
